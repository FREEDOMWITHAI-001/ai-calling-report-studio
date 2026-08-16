"""Stakeholder-facing PowerPoint: KPI slides + charts."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from . import style as S

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _rgb(hex_code: str) -> RGBColor:
    return RGBColor.from_string(hex_code)


INK = _rgb(S.INK)
MUTED = _rgb(S.MUTED)
ACCENT = _rgb(S.ACCENT)
BAND = _rgb(S.BAND)
WHITE = _rgb("FFFFFF")


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _text(slide, left, top, width, height, text, size=18, bold=False, color=INK,
          align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = wrap
    para = frame.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return box


def _card(slide, left, top, width, height, label, value, sub=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = BAND
    shape.line.color.rgb = _rgb(S.RULE)
    shape.shadow.inherit = False
    frame = shape.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.22)
    frame.margin_top = Inches(0.16)

    p0 = frame.paragraphs[0]
    r0 = p0.add_run()
    r0.text = label.upper()
    r0.font.size = Pt(10)
    r0.font.color.rgb = MUTED
    r0.font.name = "Segoe UI"

    p1 = frame.add_paragraph()
    r1 = p1.add_run()
    r1.text = value
    r1.font.size = Pt(30)
    r1.font.bold = True
    r1.font.color.rgb = INK
    r1.font.name = "Segoe UI"

    if sub:
        p2 = frame.add_paragraph()
        r2 = p2.add_run()
        r2.text = sub
        r2.font.size = Pt(10)
        r2.font.color.rgb = MUTED
        r2.font.name = "Segoe UI"
    return shape


def _title_slide(prs, result):
    slide = _blank(prs)
    meta = result["meta"]
    head = result["headline"]
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(2.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = INK
    bar.line.fill.background()
    bar.shadow.inherit = False
    _text(slide, Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.8),
          "WHAT AI CALLING ADDED", 34, True, WHITE)
    _text(slide, Inches(0.8), Inches(1.35), Inches(11.7), Inches(0.6),
          f"{meta['title']}  ·  {S.pretty_date(meta['date_from'])} – {S.pretty_date(meta['date_to'])}",
          15, False, _rgb(S.ACCENT_SOFT))

    cards = [
        ("Revenue with AI", S.money(head["revenue_with_ai"]), f"{head['buyers']:,} buyers"),
        ("Revenue without AI", S.money(head["revenue_without_ai"]), "counterfactual"),
        ("AI calling added", S.money(head["revenue_added"]), f"{head['extra_sales']:.1f} extra sales"),
        ("ROI", S.multiple(head["roi"]), f"talk cost {S.money(head['talk_cost'])}"),
    ]
    left = Inches(0.8)
    width = Inches(2.85)
    for label, value, sub in cards:
        _card(slide, left, Inches(2.95), width, Inches(1.55), label, value, sub)
        left += width + Inches(0.23)

    _text(slide, Inches(0.8), Inches(4.85), Inches(11.7), Inches(1.6), S.subtitle_line(result), 12, False, MUTED)
    return slide


def _chart_slide(prs, result):
    slide = _blank(prs)
    _text(slide, Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.6),
          "Show-up and buy rate by group", 26, True)
    groups = result["groups"]
    order = ["signup", "day_of", "both", "baseline"]
    labels = [groups[key]["label"].split("  ")[0] for key in order]

    data = CategoryChartData()
    data.categories = labels
    data.add_series("Show-up %", tuple(groups[k]["show_rate"] * 100 for k in order))
    data.add_series("Buy %", tuple(groups[k]["buy_rate"] * 100 for k in order))
    frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.7), Inches(1.2), Inches(11.9), Inches(4.6), data
    )
    chart = frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = _rgb("93C5FD")
    chart.series[1].format.fill.solid()
    chart.series[1].format.fill.fore_color.rgb = ACCENT

    sig = result["significance"]
    _text(slide, Inches(0.7), Inches(6.0), Inches(11.9), Inches(1.0),
          f"Connected vs baseline — show-up {S.p_value(sig['show_up']['p_value'])}, "
          f"buying {S.p_value(sig['buying']['p_value'])} (two-proportion z-test). "
          f"Bot groups overlap: a lead reached by both appears in each.", 11, False, MUTED)
    return slide


def _daily_chart_slide(prs, result):
    slide = _blank(prs)
    _text(slide, Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.6),
          "Registrants per day — connected vs baseline", 26, True)
    days = result["daily"]
    categories, connected, baseline = [], [], []
    for day in days:
        rows = {row.get("key"): row for row in day["rows"]}
        total_conn = 0
        for key in ("signup", "day_of"):
            if key in rows:
                total_conn += rows[key]["registrants"]
        both = rows.get("both", {}).get("registrants", 0)
        categories.append(S.pretty_date(day["date"])[:6])
        connected.append(total_conn - both)
        baseline.append(rows.get("baseline", {}).get("registrants", 0))

    data = CategoryChartData()
    data.categories = categories
    data.add_series("Connected", tuple(connected))
    data.add_series("Baseline", tuple(baseline))
    frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED, Inches(0.7), Inches(1.2), Inches(11.9), Inches(5.2), data
    )
    chart = frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = ACCENT
    chart.series[1].format.fill.solid()
    chart.series[1].format.fill.fore_color.rgb = _rgb("CBD5E1")
    return slide


def _table_slide(prs, result):
    slide = _blank(prs)
    _text(slide, Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.6),
          "Group detail (whole window)", 26, True)
    groups = result["groups"]
    order = ["total", "signup", "day_of", "both", "baseline"]
    headers = ["Group", "Registrants", "Showed", "Show-up %", "Buyers", "Buyer %", "Buyer Δ"]
    rows = len(order) + 1
    table_shape = slide.shapes.add_table(rows, len(headers), Inches(0.7), Inches(1.25),
                                         Inches(11.9), Inches(0.5 + 0.45 * rows))
    table = table_shape.table
    for idx, header in enumerate(headers):
        cell = table.cell(0, idx)
        cell.text = header
        para = cell.text_frame.paragraphs[0]
        para.runs[0].font.size = Pt(12)
        para.runs[0].font.bold = True

    for r, key in enumerate(order, start=1):
        item = groups[key]
        values = [
            item["label"],
            f"{item['registrants']:,}",
            f"{item['showed']:,}",
            S.pct(item["show_rate"]),
            f"{item['buyers']:,}",
            S.pct(item["buy_rate"], 2),
            S.delta(item["buy_delta"]),
        ]
        for c, value in enumerate(values):
            cell = table.cell(r, c)
            cell.text = value
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(11)
            run.font.bold = key == "total"
    return slide


def _method_slide(prs, result):
    slide = _blank(prs)
    head = result["headline"]
    calls = result["calls"]
    _text(slide, Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.6), "How this is calculated", 26, True)
    lines = [
        S.method_note(result),
        "",
        f"Calls: {calls['calls_placed']:,} placed · {calls['calls_connected']:,} connected · "
        f"{calls['billed_minutes']:,} billed minutes × {S.money(calls['cost_per_minute'], 2)} = "
        f"{S.money(calls['talk_cost'])}.",
        f"Break-even is {head['breakeven_sales']:.1f} sales; the AI is credited with {head['extra_sales']:.1f}.",
        "",
        "Caveat: observational, not an experiment — connected leads answered their phone, which already marks "
        "them as more engaged. A hold-out test (leave ~15% of registrants uncalled for a month) would settle it.",
    ]
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(11.9), Inches(5.4))
    frame = box.text_frame
    frame.word_wrap = True
    for idx, line in enumerate(lines):
        para = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        run = para.add_run()
        run.text = line
        run.font.size = Pt(13)
        run.font.color.rgb = INK if idx in (0,) else MUTED
        run.font.name = "Segoe UI"
        para.space_after = Pt(8)
    return slide


def build_pptx(result: dict, path: str | Path) -> Path:
    path = Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    _title_slide(prs, result)
    _chart_slide(prs, result)
    _table_slide(prs, result)
    if result["daily"]:
        _daily_chart_slide(prs, result)
    _method_slide(prs, result)
    prs.save(str(path))
    return path
